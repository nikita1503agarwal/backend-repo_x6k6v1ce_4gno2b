from fastapi import FastAPI, HTTPException, Depends, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from typing import List, Optional
from datetime import datetime, timedelta
import io
import secrets
from passlib.context import CryptContext
from jose import jwt
from pydantic import BaseModel

from database import create_document, get_documents, get_document, update_document, delete_document
from schemas import User, Project, MediaAsset, ShareLink, AuthPayload, SlideExportRequest

SECRET_KEY = "supersecretkey"  # for demo
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

app = FastAPI(title="Event Storyboard API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class Token(BaseModel):
    access_token: str
    token_type: str = "bearer"

# Simple in-API token creation for demo; in production use OAuth flows

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    expire = datetime.utcnow() + (expires_delta or timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

@app.get("/test")
async def test():
    return {"status": "ok"}

# Auth Endpoints (email + password demo; Google SSO placeholder token exchange)
@app.post("/auth/register", response_model=Token)
async def register(payload: AuthPayload):
    existing = await get_documents("user", {"email": payload.email}, 1)
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    hashed = pwd_context.hash(payload.password or secrets.token_hex(8))
    user = User(email=payload.email, name=payload.name, hashed_password=hashed)
    doc = await create_document("user", user.dict())
    token = create_access_token({"sub": doc["id"], "email": doc["email"]})
    return Token(access_token=token)

@app.post("/auth/login", response_model=Token)
async def login(payload: AuthPayload):
    users = await get_documents("user", {"email": payload.email}, 1)
    if not users:
        raise HTTPException(status_code=400, detail="Invalid credentials")
    user = users[0]
    if not payload.password or not pwd_context.verify(payload.password, user.get("hashed_password", "")):
        raise HTTPException(status_code=400, detail="Invalid credentials")
    token = create_access_token({"sub": user["id"], "email": user["email"]})
    return Token(access_token=token)

@app.post("/auth/google", response_model=Token)
async def google_sso(id_token: str = Form(...)):
    # In production, verify id_token with Google; here we accept and create/fetch user by email claim placeholder
    # For demo, we treat id_token as email
    email = id_token
    users = await get_documents("user", {"email": email}, 1)
    if users:
        user = users[0]
    else:
        user = (await create_document("user", User(email=email, provider="google").dict()))
    token = create_access_token({"sub": user["id"], "email": user["email"]})
    return Token(access_token=token)

# Projects CRUD
@app.post("/projects", response_model=Project)
async def create_project(project: Project):
    now = datetime.utcnow()
    project.created_at = now
    project.updated_at = now
    doc = await create_document("project", project.dict())
    return Project(**doc)

@app.get("/projects", response_model=List[Project])
async def list_projects(owner_id: Optional[str] = None):
    projects = await get_documents("project", {"owner_id": owner_id} if owner_id else {})
    return [Project(**p) for p in projects]

@app.get("/projects/{project_id}", response_model=Project)
async def get_project(project_id: str):
    proj = await get_document("project", {"id": project_id})
    if not proj:
        raise HTTPException(status_code=404, detail="Project not found")
    return Project(**proj)

@app.put("/projects/{project_id}", response_model=Project)
async def update_project(project_id: str, project: Project):
    project.updated_at = datetime.utcnow()
    updated = await update_document("project", {"id": project_id}, project.dict())
    if not updated:
        raise HTTPException(status_code=404, detail="Project not found")
    return Project(**updated)

@app.delete("/projects/{project_id}")
async def delete_project(project_id: str):
    ok = await delete_document("project", {"id": project_id})
    return {"success": ok}

# Media upload (store to local for demo) and metadata record
@app.post("/media/upload", response_model=MediaAsset)
async def upload_media(owner_id: str = Form(...), project_id: Optional[str] = Form(None), file: UploadFile = File(...)):
    content = await file.read()
    fname = f"uploads_{secrets.token_hex(4)}_{file.filename}"
    with open(fname, "wb") as f:
        f.write(content)
    asset = MediaAsset(owner_id=owner_id, project_id=project_id, url=fname, type='video' if file.content_type.startswith('video') else 'image', name=file.filename)
    doc = await create_document("mediaasset", asset.dict())
    return MediaAsset(**doc)

@app.get("/media", response_model=List[MediaAsset])
async def list_media(owner_id: Optional[str] = None, project_id: Optional[str] = None):
    filt = {}
    if owner_id:
        filt["owner_id"] = owner_id
    if project_id:
        filt["project_id"] = project_id
    assets = await get_documents("mediaasset", filt)
    return [MediaAsset(**a) for a in assets]

# Share Permissions
@app.post("/share", response_model=ShareLink)
async def create_share_link(project_id: str = Form(...), role: str = Form('viewer')):
    token = secrets.token_urlsafe(16)
    link = ShareLink(project_id=project_id, token=token, role=role, created_at=datetime.utcnow(), expires_at=datetime.utcnow() + timedelta(days=14))
    doc = await create_document("sharelink", link.dict())
    return ShareLink(**doc)

@app.get("/share/{token}", response_model=Project)
async def get_shared_project(token: str):
    links = await get_documents("sharelink", {"token": token}, 1)
    if not links:
        raise HTTPException(status_code=404, detail="Link not found")
    link = links[0]
    proj = await get_document("project", {"id": link["project_id"]})
    if not proj:
        raise HTTPException(status_code=404, detail="Project not found")
    return Project(**proj)

# Export endpoints (basic demo implementations)
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

@app.post("/export")
async def export_storyboard(req: SlideExportRequest):
    proj = await get_document("project", {"id": req.project_id})
    if not proj:
        raise HTTPException(status_code=404, detail="Project not found")
    slides: List[dict] = proj.get("slides", [])

    if req.format == 'images':
        # Create a simple image for each slide
        zip_bytes = io.BytesIO()
        from zipfile import ZipFile, ZIP_DEFLATED
        with ZipFile(zip_bytes, 'w', ZIP_DEFLATED) as zipf:
            for idx, s in enumerate(slides or [{}], start=1):
                img = Image.new('RGB', (1080, 1920), color=s.get('bg', '#111827'))
                draw = ImageDraw.Draw(img)
                text = s.get('text', f'Slide {idx}')
                draw.text((50, 50), text, fill=s.get('color', '#ffffff'))
                buf = io.BytesIO()
                img.save(buf, format='PNG')
                buf.seek(0)
                zipf.writestr(f'slide_{idx}.png', buf.read())
        zip_bytes.seek(0)
        return StreamingResponse(zip_bytes, media_type='application/zip', headers={'Content-Disposition': 'attachment; filename="slides.zip"'})

    if req.format == 'pptx':
        prs = Presentation()
        for idx, s in enumerate(slides or [{}], start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            tf = txBox.text_frame
            tf.text = s.get('text', f'Slide {idx}')
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return StreamingResponse(buf, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': 'attachment; filename="storyboard.pptx"'})

    # For video, return not implemented demo
    if req.format == 'video':
        raise HTTPException(status_code=501, detail="Video export not implemented in demo")

    raise HTTPException(status_code=400, detail="Invalid format")
